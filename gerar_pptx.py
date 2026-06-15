import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paleta (tema dark do projeto) ──────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x11, 0x17)
SURFACE  = RGBColor(0x16, 0x1B, 0x22)
SURFACE2 = RGBColor(0x21, 0x26, 0x2D)
ACCENT   = RGBColor(0x00, 0xC8, 0xA0)
TEXT     = RGBColor(0xE6, 0xED, 0xF3)
MUTED    = RGBColor(0x8B, 0x94, 0x9E)
BORDER   = RGBColor(0x30, 0x36, 0x3D)
ACCENT2  = RGBColor(0x00, 0x99, 0x80)
ACCENT3  = RGBColor(0x00, 0x6E, 0x5C)
YELLOW   = RGBColor(0xFF, 0xD7, 0x6E)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Helpers ────────────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill=None, lc=None, lw=Pt(0.75)):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if lc:
        sh.line.color.rgb = lc
        sh.line.width = lw
    else:
        sh.line.fill.background()
    return sh

def textbox(slide, text, x, y, w, h,
            size=Pt(13), bold=False, color=TEXT,
            align=PP_ALIGN.LEFT, italic=False, mono=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Consolas" if mono else "Inter"
    return tb

def add_para(tf, text, size=Pt(13), bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, sp=Pt(2), italic=False, mono=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = sp
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Consolas" if mono else "Inter"
    return p

def slide_base(slide):
    rect(slide, 0, 0, W, H, fill=BG)
    rect(slide, 0, 0, Inches(0.12), H, fill=ACCENT)

def header(slide, title, subtitle=None, time_label=None):
    rect(slide, Inches(0.12), 0, W - Inches(0.12), Inches(1.55), fill=SURFACE)
    rect(slide, Inches(0.12), Inches(1.55), W - Inches(0.12), Inches(0.025), fill=BORDER)
    textbox(slide, title,
            Inches(0.45), Inches(0.2), Inches(10.5), Inches(0.7),
            size=Pt(27), bold=True, color=TEXT)
    if subtitle:
        textbox(slide, subtitle,
                Inches(0.45), Inches(0.87), Inches(10), Inches(0.45),
                size=Pt(14), color=ACCENT)
    if time_label:
        textbox(slide, time_label,
                Inches(10.8), Inches(0.58), Inches(2.3), Inches(0.36),
                size=Pt(11), color=MUTED, italic=True, align=PP_ALIGN.RIGHT)

def card(slide, x, y, w, h, title=None, items=None,
         title_color=ACCENT, item_color=TEXT, bg=SURFACE,
         item_size=Pt(12), title_size=Pt(13), icon="•"):
    rect(slide, x, y, w, h, fill=bg, lc=BORDER)
    tb = slide.shapes.add_textbox(
        x + Inches(0.2), y + Inches(0.15),
        w - Inches(0.38), h - Inches(0.25))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.space_before = Pt(0)
        r = p.add_run()
        r.text = title
        r.font.size = title_size
        r.font.bold = True
        r.font.color.rgb = title_color
        r.font.name = "Inter"
        first = False
    if items:
        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(5)
            r = p.add_run()
            r.text = f"{icon}  {item}"
            r.font.size = item_size
            r.font.color.rgb = item_color
            r.font.name = "Inter"

def code_block(slide, x, y, w, h, lines, label=""):
    rect(slide, x, y, w, h, fill=SURFACE2, lc=BORDER)
    rect(slide, x, y, w, Inches(0.3), fill=BG, lc=BORDER)
    for i, (dc) in enumerate([RGBColor(255,95,87), RGBColor(255,189,46), RGBColor(40,202,65)]):
        d = slide.shapes.add_shape(9,
            x + Inches(0.16 + i * 0.2), y + Inches(0.08),
            Inches(0.12), Inches(0.12))
        d.fill.solid(); d.fill.fore_color.rgb = dc; d.line.fill.background()
    if label:
        textbox(slide, label,
                x + Inches(0.7), y + Inches(0.04), w - Inches(0.8), Inches(0.24),
                size=Pt(8.5), color=MUTED, mono=True, align=PP_ALIGN.CENTER)
    tb = slide.shapes.add_textbox(
        x + Inches(0.2), y + Inches(0.34),
        w - Inches(0.35), h - Inches(0.46))
    tf = tb.text_frame
    tf.word_wrap = False
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(0.5)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(10.5)
        r.font.name = "Consolas"
        s = line.lstrip()
        if s.startswith("//") or s.startswith("#"):
            r.font.color.rgb = MUTED
        elif s.startswith(("const ", "let ", "function ", "return ",
                            "if ", "while ", "async ", "import ", "export ",
                            "class ", "static ")):
            r.font.color.rgb = ACCENT
        elif s.startswith(("<", "</")):
            r.font.color.rgb = YELLOW
        else:
            r.font.color.rgb = TEXT

def page_num(slide, n, total=6):
    textbox(slide, f"{n} / {total}",
            Inches(12.5), Inches(7.15), Inches(0.7), Inches(0.28),
            size=Pt(9), color=MUTED, align=PP_ALIGN.RIGHT)

# ==============================================================================
# SLIDE 1 — Apresentação do Projeto
# ==============================================================================
s1 = prs.slides.add_slide(BLANK)
slide_base(s1)
header(s1,
    "EducaInclui Pascal",
    "Plataforma de cursos online desenvolvida como projeto universitario",
    "~1:30 min")
page_num(s1, 1)

# O que e o projeto
card(s1, Inches(0.3), Inches(1.7), Inches(6.2), Inches(2.65),
    title="O que e o projeto?",
    items=[
        "Site de cursos online com video aulas interativas",
        "Administradores criam cursos, aulas e perguntas",
        "Alunos assistem aos videos e respondem perguntas",
        "Perguntas aparecem automaticamente conforme o video avanca",
        "Anotacoes salvas por aula, progresso registrado localmente",
    ])

# Paginas do sistema
card(s1, Inches(0.3), Inches(4.5), Inches(6.2), Inches(2.85),
    title="Paginas do Sistema",
    items=[
        "Landing Page  —  apresentacao + login",
        "Dashboard Admin  —  resumo dos cursos cadastrados",
        "Gerenciar Cursos  —  criar, editar e excluir cursos/aulas",
        "Area do Aluno  —  buscar e filtrar cursos disponiveis",
        "Player de Video  —  assistir aulas e responder perguntas",
    ])

# Tecnologias
rect(s1, Inches(6.65), Inches(1.7), Inches(6.45), Inches(5.65), fill=SURFACE, lc=BORDER)
textbox(s1, "Tecnologias Utilizadas",
        Inches(6.85), Inches(1.85), Inches(6.0), Inches(0.38),
        size=Pt(14), bold=True, color=ACCENT)

tech_rows = [
    ("Vite 5",          "Ferramenta de build — servidor de dev rapido"),
    ("React 18",        "Biblioteca para construir a interface"),
    ("React Router v6", "Controla a navegacao entre paginas"),
    ("localStorage",    "Salva dados do usuario no navegador"),
    ("IndexedDB",       "Armazena videos e imagens localmente"),
    ("CSS puro",        "Estilizacao com Grid, Flexbox e variaveis CSS"),
    ("Vercel",          "Hospedagem gratuita na nuvem"),
]
for i, (tech, desc) in enumerate(tech_rows):
    ry = Inches(2.38 + i * 0.66)
    if i % 2 == 0:
        rect(s1, Inches(6.72), ry, Inches(6.2), Inches(0.56), fill=SURFACE2, lc=BORDER, lw=Pt(0.5))
    textbox(s1, tech,  Inches(6.85), ry + Inches(0.12), Inches(2.1), Inches(0.36),
            size=Pt(11.5), bold=True, color=ACCENT)
    textbox(s1, desc,  Inches(9.1),  ry + Inches(0.12), Inches(3.6), Inches(0.36),
            size=Pt(11), color=TEXT)

# ==============================================================================
# SLIDE 2 — Estrutura de Dados: Fila (Queue)
# ==============================================================================
s2 = prs.slides.add_slide(BLANK)
slide_base(s2)
header(s2,
    "Estrutura de Dados: Fila (Queue - FIFO)",
    "Como as perguntas sao liberadas conforme o video avanca",
    "~1:45 min")
page_num(s2, 2)

card(s2, Inches(0.3), Inches(1.7), Inches(4.1), Inches(1.55),
    title="O Problema",
    items=[
        "Cada aula pode ter varias perguntas",
        "Elas nao podem aparecer todas de uma vez",
        "Devem surgir em ordem conforme o video passa",
        "O video pausa ate o aluno responder",
    ], item_size=Pt(12))

card(s2, Inches(0.3), Inches(3.38), Inches(4.1), Inches(1.55),
    title="A Solucao: Fila FIFO",
    items=[
        "FIFO = First In, First Out (primeiro a entrar, primeiro a sair)",
        "As perguntas entram na fila ao carregar a aula",
        "O video monitora seu progresso (evento timeupdate)",
        "Ao passar do ponto certo, a 1a pergunta da fila e exibida",
    ], item_size=Pt(12))

card(s2, Inches(0.3), Inches(5.07), Inches(4.1), Inches(2.28),
    title="Como funciona na pratica",
    items=[
        "Aula com 3 perguntas: fila = [P1, P2, P3]",
        "Video em 20%: nenhuma pergunta ainda",
        "Video em 35%: P1 sai da fila -> video pausa -> aluno responde",
        "Video retoma -> em 55%: P2 aparece",
        "Video em 75%: P3 aparece -> ultima pergunta",
        "Fila vazia -> video segue ate o fim",
    ], item_size=Pt(11.5))

code_block(s2, Inches(4.6), Inches(1.7), Inches(8.5), Inches(5.65),
    label="Queue.js + CourseViewPage.jsx",
    lines=[
        "// Fila FIFO — estrutura de dados do projeto",
        "class Queue {",
        "  #items = []",
        "",
        "  enqueue(item) { this.#items.push(item) }   // adiciona no fim",
        "  dequeue()     { return this.#items.shift() } // remove do inicio",
        "  isEmpty()     { return this.#items.length === 0 }",
        "  get size()    { return this.#items.length }",
        "}",
        "",
        "// Ao carregar uma aula, todas as perguntas entram na fila",
        "lesson.questions.forEach(q => questionQueue.enqueue(q))",
        "",
        "// O video avanca e monitora o progresso",
        "video.addEventListener('timeupdate', () => {",
        "  const progresso = video.currentTime / video.duration",
        "",
        "  // Chegou no ponto da proxima pergunta?",
        "  if (progresso >= thresholds[unlockedCount]) {",
        "    const pergunta = questionQueue.dequeue() // tira da fila",
        "    video.pause()                            // pausa o video",
        "    mostrarOverlay(pergunta)                 // exibe a pergunta",
        "    unlockedCount++",
        "  }",
        "})",
    ])

# ==============================================================================
# SLIDE 3 — React: Componentes e Hooks
# ==============================================================================
s3 = prs.slides.add_slide(BLANK)
slide_base(s3)
header(s3,
    "React: Componentes e Gerenciamento de Estado",
    "Como a interface e organizada e como os dados fluem",
    "~1:45 min")
page_num(s3, 3)

card(s3, Inches(0.3), Inches(1.7), Inches(4.15), Inches(1.75),
    title="O que sao Componentes?",
    items=[
        "Blocos reutilizaveis de interface (como 'pecas de Lego')",
        "Cada componente tem sua propria logica e visual",
        "Podem ser combinados para formar paginas completas",
        "Ex: QuestionItem, LessonModal, AddCourseModal",
    ], item_size=Pt(12))

card(s3, Inches(0.3), Inches(3.6), Inches(4.15), Inches(1.62),
    title="Organizacao do Projeto",
    items=[
        "Pages  ->  telas principais (5 rotas)",
        "Modals ->  janelas de adicionar/editar",
        "Items  ->  partes menores reutilizaveis",
        "Libs   ->  servicos de dados (localStorage, IndexedDB)",
    ], item_size=Pt(12))

# Tabela de hooks
rect(s3, Inches(0.3), Inches(5.35), Inches(4.15), Inches(2.0), fill=SURFACE, lc=BORDER)
textbox(s3, "Hooks Utilizados",
        Inches(0.5), Inches(5.5), Inches(3.8), Inches(0.3),
        size=Pt(13), bold=True, color=ACCENT)

hooks = [
    ("useState",   "Guarda valores que mudam a tela (modal aberto, filtro)"),
    ("useEffect",  "Executa codigo quando algo muda (carregar video, notas)"),
    ("useRef",     "Acessa o elemento <video> diretamente no HTML"),
    ("useCallback","Evita recriar funcoes desnecessariamente"),
]
for i, (h, d) in enumerate(hooks):
    hy = Inches(5.88 + i * 0.34)
    textbox(s3, h, Inches(0.5),  hy, Inches(1.35), Inches(0.3),
            size=Pt(10.5), bold=True, color=ACCENT, mono=True)
    textbox(s3, d, Inches(1.95), hy, Inches(2.35), Inches(0.3),
            size=Pt(10), color=TEXT)

code_block(s3, Inches(4.65), Inches(1.7), Inches(8.45), Inches(5.65),
    label="Exemplo: CoursesPage.jsx — adicionar curso com estado React",
    lines=[
        "// Estado da pagina — React controla o que aparece na tela",
        "const [cursos, setCursos]         = useState([])",
        "const [modalAberto, setModal]     = useState(false)",
        "const [nomeCurso, setNome]        = useState('')",
        "const [topico, setTopico]         = useState('')",
        "",
        "// Carrega os cursos do localStorage ao abrir a pagina",
        "useEffect(() => {",
        "  setCursos(CoursesDB.load())",
        "}, [])",
        "",
        "// Salva o novo curso",
        "function adicionarCurso(e) {",
        "  e.preventDefault()",
        "  const novo = new Course({ id: Date.now(), name: nomeCurso, topic: topico })",
        "  const atualizados = [...cursos, novo]",
        "  CoursesDB.save(atualizados)   // salva no localStorage",
        "  setCursos(atualizados)        // atualiza a tela",
        "  setModal(false)              // fecha o modal",
        "}",
        "",
        "// JSX — React renderiza a lista automaticamente",
        "return (",
        "  <div>",
        "    {cursos.map(c => <div key={c.id}>{c.name}</div>)}",
        "  </div>",
        ")",
    ])

# ==============================================================================
# SLIDE 4 — Arquitetura SPA e Rotas
# ==============================================================================
s4 = prs.slides.add_slide(BLANK)
slide_base(s4)
header(s4,
    "Arquitetura SPA e Navegacao",
    "Como o site funciona sem recarregar a pagina",
    "~1:30 min")
page_num(s4, 4)

card(s4, Inches(0.3), Inches(1.7), Inches(4.15), Inches(1.55),
    title="O que e uma SPA?",
    items=[
        "SPA = Single Page Application",
        "O navegador carrega o HTML uma unica vez",
        "A navegacao entre paginas e feita pelo JavaScript",
        "Mais rapido — nao precisa recarregar tudo do servidor",
    ], item_size=Pt(12))

card(s4, Inches(0.3), Inches(3.38), Inches(4.15), Inches(1.55),
    title="Protecao de Rotas",
    items=[
        "Nem toda pagina pode ser acessada por qualquer usuario",
        "Alunos (cliente) -> /user e /course/:id",
        "Admins (professor/educacional) -> /admin e /admin/courses",
        "Quem nao esta logado e redirecionado para a home",
    ], item_size=Pt(12))

card(s4, Inches(0.3), Inches(5.07), Inches(4.15), Inches(2.28),
    title="Deploy no Vercel",
    items=[
        "Problema: acessar /admin diretamente daria erro 404",
        "Em SPA, so o index.html existe no servidor",
        "Solucao: vercel.json redireciona TUDO para o index.html",
        "O React Router assume e mostra a pagina certa",
        'vercel.json: { "rewrites": [{ "source": "/(.*)", ',
        '               "destination": "/index.html" }] }',
    ], item_size=Pt(11.5))

code_block(s4, Inches(4.65), Inches(1.7), Inches(8.45), Inches(5.65),
    label="App.jsx — Rotas e Protecao de Acesso",
    lines=[
        "// Componente que protege rotas privadas",
        "function RequireAuth({ children, roles }) {",
        "  const user = Session.get()  // pega o usuario logado",
        "",
        "  if (!user)  // nao esta logado? volta para a home",
        "    return <Navigate to='/' replace />",
        "",
        "  if (!roles.includes(user.tipo))  // tipo errado? redireciona",
        "    return <Navigate to='/user' replace />",
        "",
        "  return children  // tudo certo, mostra a pagina",
        "}",
        "",
        "// Mapa de rotas da aplicacao",
        "function App() {",
        "  return (",
        "    <BrowserRouter>",
        "      <Routes>",
        "        <Route path='/'     element={<LandingPage />} />",
        "",
        "        <Route path='/admin'",
        "          element={<RequireAuth roles={['professor','educacional']}>",
        "            <AdministrativePage />",
        "          </RequireAuth>} />",
        "",
        "        <Route path='/user'",
        "          element={<RequireAuth roles={['cliente']}>",
        "            <UserPage />",
        "          </RequireAuth>} />",
        "",
        "        <Route path='/course/:courseId'",
        "          element={<RequireAuth roles={['cliente']}>",
        "            <CourseViewPage />",
        "          </RequireAuth>} />",
        "      </Routes>",
        "    </BrowserRouter>",
        "  )",
        "}",
    ])

# ==============================================================================
# SLIDE 5 — Armazenamento de Dados
# ==============================================================================
s5 = prs.slides.add_slide(BLANK)
slide_base(s5)
header(s5,
    "Armazenamento de Dados no Navegador",
    "Como o projeto guarda dados sem precisar de servidor",
    "~1:45 min")
page_num(s5, 5)

card(s5, Inches(0.3), Inches(1.7), Inches(3.95), Inches(2.0),
    title="localStorage",
    items=[
        "Guarda textos simples no navegador",
        "Persiste mesmo fechando o navegador",
        "Usado para: lista de cursos, sessao do usuario",
        "notas por aula, progresso de aulas concluidas",
        "Limite: ~5 MB por site",
    ], item_size=Pt(12))

card(s5, Inches(0.3), Inches(3.85), Inches(3.95), Inches(1.85),
    title="IndexedDB",
    items=[
        "Banco de dados real dentro do navegador",
        "Suporta arquivos grandes (videos e imagens)",
        "Usado para: armazenar arquivos .mp4 das aulas",
        "e fotos de capa dos cursos enviadas pelo admin",
        "Limite: varios GB (depende do disco)",
    ], item_size=Pt(12))

card(s5, Inches(0.3), Inches(5.85), Inches(3.95), Inches(1.5),
    title="Modelos de Dados",
    items=[
        "Course  ->  id, nome, topico, lista de aulas",
        "Lesson  ->  id, nome, descricao, video, perguntas",
        "Question ->  texto, tipo (aberta/fechada), gabarito",
        "fromRaw() reconstroi objetos salvos no localStorage",
    ], item_size=Pt(11.5))

code_block(s5, Inches(4.45), Inches(1.7), Inches(8.65), Inches(5.65),
    label="CoursesDB.js e VideoDB.js — como os dados sao salvos e lidos",
    lines=[
        "// ── localStorage: cursos e aulas ────────────────────────────",
        "const CoursesDB = {",
        "",
        "  // Le todos os cursos salvos",
        "  load() {",
        "    const raw = JSON.parse(localStorage.getItem('pascal_courses') || '[]')",
        "    return raw.map(c => Course.fromRaw(c))  // reconstroi os objetos",
        "  },",
        "",
        "  // Salva a lista inteira de cursos",
        "  save(courses) {",
        "    localStorage.setItem('pascal_courses', JSON.stringify(courses))",
        "  }",
        "}",
        "",
        "// ── IndexedDB: arquivos de video e imagem ─────────────────────",
        "const VideoDB = {",
        "",
        "  // Salva um arquivo (Blob) com uma chave",
        "  async save(chave, arquivo) {",
        "    const db = await abrirBanco()",
        "    db.transaction('videos', 'readwrite')",
        "      .objectStore('videos').put(arquivo, chave)",
        "  },",
        "",
        "  // Recupera o arquivo pela chave",
        "  async get(chave) {",
        "    const db = await abrirBanco()",
        "    // retorna o Blob que vira uma URL para o <video>",
        "    return db.transaction('videos', 'readonly')",
        "      .objectStore('videos').get(chave)",
        "  }",
        "}",
    ])

# ==============================================================================
# SLIDE 6 — Ferramentas e Resultado Final
# ==============================================================================
s6 = prs.slides.add_slide(BLANK)
slide_base(s6)
header(s6,
    "Ferramentas e Resultado Final",
    "Build, deploy e o que foi entregue",
    "~1:45 min")
page_num(s6, 6)

card(s6, Inches(0.3), Inches(1.7), Inches(3.95), Inches(1.72),
    title="Vite — Ferramenta de Build",
    items=[
        "Transforma JSX em HTML/CSS/JS que o browser entende",
        "Servidor local com atualizacao instantanea (HMR)",
        "npm run dev   -> abre o site localmente",
        "npm run build -> gera a pasta dist/ para producao",
    ], item_size=Pt(12))

card(s6, Inches(0.3), Inches(3.57), Inches(3.95), Inches(1.72),
    title="Deploy no Vercel",
    items=[
        "Hospedagem gratuita direto do GitHub",
        "Cada push na main atualiza o site automaticamente",
        "Certificado HTTPS gratuito e dominio personalizado",
        "Tempo de build: menos de 10 segundos",
    ], item_size=Pt(12))

# Resultado final
rect(s6, Inches(0.3), Inches(5.44), Inches(3.95), Inches(1.92), fill=SURFACE, lc=BORDER)
textbox(s6, "Resultado do Build de Producao",
        Inches(0.5), Inches(5.58), Inches(3.6), Inches(0.3),
        size=Pt(13), bold=True, color=ACCENT)

resultados = [
    ("JavaScript total:", "203 kB  (63 kB comprimido)"),
    ("CSS total:",        " 34 kB  ( 6 kB comprimido)"),
    ("Tempo de build:",   "  8 segundos"),
    ("Arquivos gerados:", "  3 arquivos (index.html + JS + CSS)"),
]
for i, (label, value) in enumerate(resultados):
    ry = Inches(5.98 + i * 0.31)
    textbox(s6, label, Inches(0.5),  ry, Inches(1.9), Inches(0.28),
            size=Pt(11), color=MUTED)
    textbox(s6, value, Inches(2.45), ry, Inches(1.65), Inches(0.28),
            size=Pt(11), bold=True, color=TEXT, mono=True)

code_block(s6, Inches(4.45), Inches(1.7), Inches(8.65), Inches(4.35),
    label="Fluxo completo: do codigo ao site publicado",
    lines=[
        "// 1. Desenvolvimento local",
        "npm run dev",
        "// -> abre http://localhost:5173",
        "// -> qualquer mudanca no codigo atualiza o browser na hora",
        "",
        "// 2. Gerar versao de producao",
        "npm run build",
        "// -> cria a pasta dist/ com tudo otimizado",
        "// -> JSX vira JavaScript puro",
        "// -> CSS e unificado em um arquivo so",
        "// -> imagens e assets sao copiados",
        "",
        "// 3. Publicar no Vercel (automatico pelo GitHub)",
        "git add . && git commit -m 'nova versao'",
        "git push origin main",
        "// -> Vercel detecta o push",
        "// -> executa npm run build automaticamente",
        "// -> publica em https://seu-projeto.vercel.app",
    ])

# Resumo final
rect(s6, Inches(4.45), Inches(6.18), Inches(8.65), Inches(1.18), fill=SURFACE2, lc=ACCENT, lw=Pt(1.5))
tb_r = s6.shapes.add_textbox(Inches(4.65), Inches(6.3), Inches(8.3), Inches(0.95))
tf_r = tb_r.text_frame
tf_r.word_wrap = True
p0 = tf_r.paragraphs[0]
p0.space_before = Pt(0)
r0 = p0.add_run()
r0.text = "O projeto funciona 100% no navegador, sem precisar de servidor ou banco de dados externo."
r0.font.size = Pt(12.5); r0.font.bold = True; r0.font.color.rgb = TEXT; r0.font.name = "Inter"
p1 = tf_r.add_paragraph()
p1.space_before = Pt(4)
r1 = p1.add_run()
r1.text = "Stack: Vite + React + React Router + localStorage + IndexedDB + Vercel"
r1.font.size = Pt(11.5); r1.font.color.rgb = ACCENT; r1.font.name = "Inter"

# ── Salvar ─────────────────────────────────────────────────────────────────────
out = r"c:\Users\Administrator\Desktop\DEV\PROJETOS\Front\EducaInclui_Pascal_Apresentacao.pptx"
prs.save(out)
print(f"Arquivo salvo: {out}")
